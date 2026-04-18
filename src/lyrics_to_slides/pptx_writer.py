"""PowerPoint generation boundaries."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.text.fonts import FontFiles
from pptx.text.layout import _rendered_size
from pptx.util import Cm, Pt

from .models import Slideshow

SLIDE_WIDTH_CM = 33.87
SLIDE_HEIGHT_CM = 19.05
BLANK_LAYOUT = 6
FONT_NAME = "Arial"
FONT_SIZE_PT = 60
FIT_WIDTH_RATIO = 0.96
FIT_HEIGHT_RATIO = 0.96
AMIN_TEXT = "-AMIN-"
AMIN_FONT_SIZE_PT = 24
AMIN_BOX_WIDTH_CM = 5
AMIN_BOX_HEIGHT_CM = 1
AMIN_RIGHT_MARGIN_CM = 0.4
AMIN_BOTTOM_MARGIN_CM = 0.2


def build_presentation() -> Presentation:
    """Create a presentation with the project's standard slide dimensions."""
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_WIDTH_CM)
    prs.slide_height = Cm(SLIDE_HEIGHT_CM)
    return prs


def _calculate_font_size(lines: list[str], width: int, height: int) -> int:
    """Return the largest font size that fits all lyric lines without wrapping."""
    font_file = FontFiles.find(FONT_NAME, False, False)
    max_width = int(width * FIT_WIDTH_RATIO)
    max_height = int(height * FIT_HEIGHT_RATIO)

    for point_size in range(FONT_SIZE_PT, 0, -1):
        line_height = _rendered_size("Ty", point_size, font_file)[1]
        if line_height * len(lines) > max_height:
            continue

        if all(_rendered_size(line or " ", point_size, font_file)[0] <= max_width for line in lines):
            return point_size

    return 1


def add_lyric_slide(prs: Presentation, text: str):
    """Add a lyric slide with the default projector-friendly styling."""
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])

    background_fill = slide.background.fill
    background_fill.solid()
    background_fill.fore_color.rgb = RGBColor(0, 0, 0)

    textbox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.word_wrap = False

    lines = text.splitlines() or [text]
    font_size = _calculate_font_size(lines, prs.slide_width, prs.slide_height)
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.CENTER

        run = paragraph.add_run()
        run.text = line
        run.font.name = FONT_NAME
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(255, 255, 255)

    return slide


def add_amin_footer(prs: Presentation, slide) -> None:
    """Add the closing footer tag to the bottom-right corner of a slide."""
    footer_width = Cm(AMIN_BOX_WIDTH_CM)
    footer_height = Cm(AMIN_BOX_HEIGHT_CM)
    footer_left = prs.slide_width - footer_width - Cm(AMIN_RIGHT_MARGIN_CM)
    footer_top = prs.slide_height - footer_height - Cm(AMIN_BOTTOM_MARGIN_CM)

    textbox = slide.shapes.add_textbox(footer_left, footer_top, footer_width, footer_height)
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.word_wrap = False
    text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.RIGHT

    run = paragraph.add_run()
    run.text = AMIN_TEXT
    run.font.name = FONT_NAME
    run.font.size = Pt(AMIN_FONT_SIZE_PT)
    run.font.color.rgb = RGBColor(255, 255, 255)


def write_presentation(output_path: Path, slides: Slideshow) -> Path:
    """Write a single PowerPoint file for one song/page.

    Implementation:
    - create a simple, readable lyric slide theme
    - add one slide per section in `slideshow`
    - save the deck to `output_path`

    Args:
        output_path (Path): The destination path for the generated `.pptx` file.
        slides (Slideshow): The ordered slides to include in the deck.

    Returns:
        Path: The saved PowerPoint file path.
    """
    prs = build_presentation()
    output_path.mkdir(parents=True, exist_ok=True)

    add_lyric_slide(prs, slides.title)      # Title slide

    for section in slides.body_sections:
        add_lyric_slide(prs, "\n".join(section.lines))

    if prs.slides:
        add_amin_footer(prs, prs.slides[-1])

    file_path = output_path / f"{slides.title}.pptx"
    prs.save(file_path)
    return file_path
